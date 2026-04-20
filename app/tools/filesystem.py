"""Phase 5: workspace-aware filesystem tools.

All tools resolve the given path and check it is inside WORKSPACE_DIR.
Paths outside the workspace raise OutsideWorkspaceError.
"""
from __future__ import annotations

import os
from pathlib import Path
from typing import Annotated

from langchain_core.tools import tool

import app.config as app_config

from sqlalchemy import select as sa_select
from app.db.engine import SyncSessionLocal
from app.db.models import WorkspaceLocation


def _get_allowed_roots() -> list[tuple[Path, bool]]:
    """Return [(resolved_path, writable), ...] for all workspace locations.

    Primary location is always first in the list.
    """
    try:
        with SyncSessionLocal() as db:
            rows = db.execute(
                sa_select(WorkspaceLocation).order_by(
                    WorkspaceLocation.is_primary.desc(),
                    WorkspaceLocation.created_at
                )
            ).scalars().all()
        if rows:
            return [(Path(os.path.realpath(r.path)), r.writable) for r in rows]
    except Exception:
        pass
    # Fallback: use app_config.WORKSPACE_DIR if DB is unavailable
    return [(Path(os.path.realpath(str(app_config.WORKSPACE_DIR))), True)]


class OutsideWorkspaceError(ValueError):
    """Raised when a path falls outside the allowed workspace directory."""


def _safe_resolve(path_str: str, require_writable: bool = False) -> Path:
    """Resolve *path_str* and verify it falls inside an allowed workspace location.

    If *require_writable* is True and the matching location is read-only,
    raises OutsideWorkspaceError.
    """
    if not path_str or not isinstance(path_str, str):
        raise OutsideWorkspaceError("Empty or invalid path.")

    if path_str.startswith("\\\\") or path_str.startswith("//"):
        raise OutsideWorkspaceError("UNC paths are not allowed.")
    if path_str.startswith("\\\\?\\") or path_str.startswith("\\\\.\\"):
        raise OutsideWorkspaceError("Device / extended-length paths are not allowed.")
    if ":" in path_str[2:]:
        raise OutsideWorkspaceError("Alternate data streams are not allowed.")

    roots = _get_allowed_roots()
    # Use primary root for relative-path resolution
    primary_root = roots[0][0]

    candidate = Path(path_str)
    if not candidate.is_absolute():
        candidate = primary_root / candidate

    try:
        resolved = Path(os.path.realpath(str(candidate)))
    except OSError as e:
        raise OutsideWorkspaceError(f"Could not resolve path '{path_str}': {e}")

    for root_path, writable in roots:
        try:
            resolved.relative_to(root_path)
        except ValueError:
            continue
        # Path is inside this root
        if require_writable and not writable:
            raise OutsideWorkspaceError(
                f"'{path_str}' is in a read-only workspace location '{root_path}'. "
                "This location is read-only and cannot be written to."
            )
        return resolved

    raise OutsideWorkspaceError(
        f"Path '{path_str}' is outside all allowed workspace locations. "
        "I can only access files inside configured workspace directories."
    )


def _read_docx(path: Path) -> str:
    """Extract text from a .docx Word document."""
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts) if parts else "(empty document)"


def _read_xlsx(path: Path) -> str:
    """Extract text from an .xlsx Excel workbook."""
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"=== Sheet: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                parts.append("\t".join(cells))
    return "\n".join(parts) if parts else "(empty workbook)"


def _read_pptx(path: Path) -> str:
    """Extract text from a .pptx PowerPoint presentation."""
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            parts.append(f"--- Slide {i} ---")
            parts.extend(slide_texts)
    return "\n".join(parts) if parts else "(empty presentation)"


def _read_pdf(path: Path) -> str:
    """Extract text from a PDF file."""
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"--- Page {i} ---\n{text.strip()}")
    return "\n\n".join(parts) if parts else "(no extractable text in PDF)"


def _read_csv(path: Path) -> str:
    """Read a CSV file and return it as plain text."""
    import csv
    rows = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append("\t".join(row))
    return "\n".join(rows) if rows else "(empty CSV)"


# Map of file extensions to reader functions
_BINARY_READERS = {
    ".docx": _read_docx,
    ".doc":  _read_docx,   # python-docx can handle old .doc via same API in most cases
    ".xlsx": _read_xlsx,
    ".xls":  _read_xlsx,
    ".pptx": _read_pptx,
    ".pdf":  _read_pdf,
    ".csv":  _read_csv,
}

# Extensions treated as plain text regardless of content
_TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".html", ".css", ".json",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".env", ".sh", ".bat",
    ".xml", ".rst", ".log", ".sql", ".r", ".java", ".c", ".cpp", ".h",
}


@tool
def read_file(path: Annotated[str, "Path to the file to read (relative to workspace or absolute)"]) -> str:
    """Read the contents of a file inside the workspace directory.

    Supports plain text files as well as common office formats:
    Word (.docx), Excel (.xlsx), PowerPoint (.pptx), PDF (.pdf), CSV (.csv).
    """
    try:
        resolved = _safe_resolve(path)
    except OutsideWorkspaceError as e:
        return str(e)

    if not resolved.exists():
        return f"Error: File '{path}' does not exist."
    if not resolved.is_file():
        return f"Error: '{path}' is not a file."

    suffix = resolved.suffix.lower()

    # Binary/structured formats
    if suffix in _BINARY_READERS:
        try:
            return _BINARY_READERS[suffix](resolved)
        except Exception as e:
            return f"Error reading {suffix} file: {e}"

    # Plain text (known extensions or unknown — try UTF-8)
    try:
        return resolved.read_text(encoding="utf-8", errors="replace")
    except OSError as e:
        return f"Error reading file: {e}"


@tool
def write_file(
    path: Annotated[str, "Path to write (relative to workspace or absolute)"],
    content: Annotated[str, "Content to write to the file"],
) -> str:
    """Write content to a file inside the workspace directory. Creates the file if it does not exist."""
    try:
        resolved = _safe_resolve(path, require_writable=True)
    except OutsideWorkspaceError as e:
        return str(e)

    try:
        resolved.parent.mkdir(parents=True, exist_ok=True)
        resolved.write_text(content, encoding="utf-8")
        return f"Successfully wrote {len(content)} characters to '{path}'."
    except OSError as e:
        return f"Error writing file: {e}"


@tool
def list_dir(path: Annotated[str, "Directory path to list (relative to workspace or absolute, default '.')"] = ".") -> str:
    """List the files and directories inside a workspace directory."""
    try:
        resolved = _safe_resolve(path)
    except OutsideWorkspaceError as e:
        return str(e)

    if not resolved.exists():
        return f"Error: Directory '{path}' does not exist."
    if not resolved.is_dir():
        return f"Error: '{path}' is not a directory."

    try:
        entries = sorted(resolved.iterdir(), key=lambda p: (p.is_file(), p.name))
        lines = []
        for entry in entries:
            kind = "FILE" if entry.is_file() else "DIR "
            size = ""
            if entry.is_file():
                try:
                    size = f"  ({entry.stat().st_size} bytes)"
                except OSError:
                    pass
            lines.append(f"[{kind}] {entry.name}{size}")
        if not lines:
            return f"Directory '{path}' is empty."
        return "\n".join(lines)
    except OSError as e:
        return f"Error listing directory: {e}"


@tool
def create_folder(path: Annotated[str, "Path of the folder to create (relative to workspace or absolute)"]) -> str:
    """Create a folder (directory) inside the workspace directory. Creates all intermediate directories as needed."""
    try:
        resolved = _safe_resolve(path, require_writable=True)
    except OutsideWorkspaceError as e:
        return str(e)

    try:
        resolved.mkdir(parents=True, exist_ok=True)
        return f"Successfully created folder '{path}'."
    except OSError as e:
        return f"Error creating folder: {e}"


@tool
def delete_file(path: Annotated[str, "Path to the file to delete (relative to workspace or absolute)"]) -> str:
    """Delete a file inside the workspace directory. Directories are not deleted by this tool."""
    try:
        resolved = _safe_resolve(path, require_writable=True)
    except OutsideWorkspaceError as e:
        return str(e)

    if not resolved.exists():
        return f"Error: File '{path}' does not exist."
    if not resolved.is_file():
        return f"Error: '{path}' is a directory, not a file. This tool only deletes files."

    try:
        resolved.unlink()
        return f"Successfully deleted '{path}'."
    except OSError as e:
        return f"Error deleting file: {e}"


@tool
def copy_file(
    src: Annotated[str, "Source file path (relative to workspace or absolute)"],
    dst: Annotated[str, "Destination path (relative to workspace or absolute). Can be a file path or a directory."],
) -> str:
    """Copy a file within the workspace. Preserves binary content exactly (safe for images, PDFs, etc.).
    If dst is a directory, the file is copied into it with its original name.
    Creates destination directories as needed."""
    import shutil
    try:
        src_resolved = _safe_resolve(src)
        dst_resolved = _safe_resolve(dst, require_writable=True)
    except OutsideWorkspaceError as e:
        return str(e)

    if not src_resolved.exists():
        return f"Error: Source file '{src}' does not exist."
    if not src_resolved.is_file():
        return f"Error: '{src}' is not a file."

    # If dst is (or looks like) a directory, copy into it with the original name
    if dst_resolved.is_dir():
        dst_resolved = dst_resolved / src_resolved.name
    else:
        dst_resolved.parent.mkdir(parents=True, exist_ok=True)

    try:
        shutil.copy2(str(src_resolved), str(dst_resolved))
        return f"Copied '{src}' → '{dst_resolved.relative_to(app_config.WORKSPACE_DIR)}'."
    except OSError as e:
        return f"Error copying file: {e}"


@tool
def move_file(
    src: Annotated[str, "Source file path (relative to workspace or absolute)"],
    dst: Annotated[str, "Destination path (relative to workspace or absolute). Can be a file path or a directory."],
) -> str:
    """Move (rename) a file within the workspace. Safe for binary files.
    If dst is a directory, the file is moved into it with its original name.
    Creates destination directories as needed."""
    import shutil
    try:
        src_resolved = _safe_resolve(src, require_writable=True)
        dst_resolved = _safe_resolve(dst, require_writable=True)
    except OutsideWorkspaceError as e:
        return str(e)

    if not src_resolved.exists():
        return f"Error: Source file '{src}' does not exist."
    if not src_resolved.is_file():
        return f"Error: '{src}' is not a file."

    if dst_resolved.is_dir():
        dst_resolved = dst_resolved / src_resolved.name
    else:
        dst_resolved.parent.mkdir(parents=True, exist_ok=True)

    try:
        shutil.move(str(src_resolved), str(dst_resolved))
        return f"Moved '{src}' → '{dst_resolved.relative_to(app_config.WORKSPACE_DIR)}'."
    except OSError as e:
        return f"Error moving file: {e}"


@tool
def find_file(
    filename: Annotated[str, "Filename to search for (e.g. 'resume.pdf'). Case-insensitive on Windows."],
) -> str:
    """Search for a file by name across all workspace locations.
    Searches the primary workspace first, then secondary locations in order added.
    Returns the full path and which workspace it was found in, or a not-found message.
    """
    roots = _get_allowed_roots()
    filename_lower = filename.lower()
    results: list[str] = []

    for root_path, _writable in roots:
        if not root_path.exists():
            continue
        for dirpath, _dirs, files in os.walk(str(root_path)):
            for fname in files:
                if fname.lower() == filename_lower:
                    full = Path(dirpath) / fname
                    results.append(str(full))

    if not results:
        return f"File '{filename}' not found in any workspace location."
    if len(results) == 1:
        return f"Found: {results[0]}"
    lines = [f"Found {len(results)} matches:"]
    for p in results:
        lines.append(f"  {p}")
    return "\n".join(lines)
